from __future__ import annotations

import os
from datetime import datetime
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from ..models.schema import Candidate, Evaluation, Interview, Job, RankedCandidate

# ─── Brand colours ────────────────────────────────────────────────────────────
DARK    = RGBColor(0x1A, 0x1A, 0x2E)   # deep navy
ACCENT  = RGBColor(0x16, 0x21, 0x3E)   # mid navy
GOLD    = RGBColor(0xE9, 0x4F, 0x37)   # vivid red-orange (accent)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT   = RGBColor(0xF5, 0xF5, 0xF5)
GREY    = RGBColor(0x88, 0x88, 0x99)


# ─── Helpers ──────────────────────────────────────────────────────────────────

def _add_slide(prs: Presentation, layout_idx: int = 6):
    """Add a blank slide (layout 6 = blank)."""
    layout = prs.slide_layouts[layout_idx]
    return prs.slides.add_slide(layout)


def _fill_bg(slide, color: RGBColor):
    from pptx.util import Emu
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_textbox(slide, text: str, left, top, width, height,
                 font_size=18, bold=False, color=WHITE,
                 align=PP_ALIGN.LEFT, italic=False):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox


def _score_bar(slide, label: str, score: float, left, top, width=Inches(3)):
    """Draw a labelled score bar."""
    from pptx.util import Emu
    bar_height = Inches(0.18)
    # Label
    _add_textbox(slide, f"{label}", left, top, width, Inches(0.25),
                 font_size=10, color=LIGHT)
    # Background track
    bg = slide.shapes.add_shape(1, left, top + Inches(0.27), width, bar_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = ACCENT
    bg.line.fill.background()
    # Score fill
    fill_w = int(width * (score / 10))
    if fill_w > 0:
        bar = slide.shapes.add_shape(1, left, top + Inches(0.27), fill_w, bar_height)
        bar.fill.solid()
        bar.fill.fore_color.rgb = GOLD
        bar.line.fill.background()
    # Score text
    _add_textbox(slide, f"{score:.1f}/10", left + width + Inches(0.1),
                 top + Inches(0.2), Inches(0.6), Inches(0.25),
                 font_size=10, bold=True, color=GOLD)


# ─── Slides ───────────────────────────────────────────────────────────────────

def _slide_cover(prs: Presentation, job: Job, client_name: str):
    slide = _add_slide(prs)
    _fill_bg(slide, DARK)
    W, H = prs.slide_width, prs.slide_height

    # Accent bar
    bar = slide.shapes.add_shape(1, 0, H * 0.6, W, Inches(0.06))
    bar.fill.solid()
    bar.fill.fore_color.rgb = GOLD
    bar.line.fill.background()

    _add_textbox(slide, "CANDIDATE INTELLIGENCE REPORT",
                 Inches(0.8), Inches(1.2), W - Inches(1.6), Inches(0.6),
                 font_size=14, color=GREY, align=PP_ALIGN.CENTER)

    _add_textbox(slide, job.title,
                 Inches(0.8), Inches(1.9), W - Inches(1.6), Inches(1.2),
                 font_size=40, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    _add_textbox(slide, job.company,
                 Inches(0.8), Inches(3.1), W - Inches(1.6), Inches(0.5),
                 font_size=22, color=GOLD, align=PP_ALIGN.CENTER)

    date_str = datetime.now().strftime("%B %Y")
    _add_textbox(slide, f"Prepared for {client_name}  ·  {date_str}",
                 Inches(0.8), H - Inches(1.2), W - Inches(1.6), Inches(0.4),
                 font_size=11, color=GREY, align=PP_ALIGN.CENTER)


def _slide_overview(prs: Presentation, job: Job, ranked: list[RankedCandidate]):
    slide = _add_slide(prs)
    _fill_bg(slide, DARK)
    W, H = prs.slide_width, prs.slide_height

    _add_textbox(slide, "PROCESS OVERVIEW", Inches(0.6), Inches(0.4),
                 W - Inches(1.2), Inches(0.4), font_size=11, color=GREY)
    _add_textbox(slide, "Hiring Summary", Inches(0.6), Inches(0.8),
                 W - Inches(1.2), Inches(0.6), font_size=28, bold=True, color=WHITE)

    stats = [
        ("CVs Reviewed", str(len(ranked))),
        ("Shortlisted", str(sum(1 for r in ranked if r.score >= 7))),
        ("Avg Score", f"{sum(r.score for r in ranked)/len(ranked):.1f}" if ranked else "–"),
        ("Top Score", f"{ranked[0].score:.1f}" if ranked else "–"),
    ]

    box_w = Inches(1.8)
    for i, (label, value) in enumerate(stats):
        x = Inches(0.6) + i * (box_w + Inches(0.3))
        y = Inches(1.8)
        box = slide.shapes.add_shape(1, x, y, box_w, Inches(1.1))
        box.fill.solid()
        box.fill.fore_color.rgb = ACCENT
        box.line.fill.background()
        _add_textbox(slide, value, x, y + Inches(0.1), box_w, Inches(0.55),
                     font_size=32, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
        _add_textbox(slide, label, x, y + Inches(0.65), box_w, Inches(0.35),
                     font_size=11, color=LIGHT, align=PP_ALIGN.CENTER)

    # Rankings table header
    _add_textbox(slide, "CANDIDATE RANKINGS", Inches(0.6), Inches(3.2),
                 W - Inches(1.2), Inches(0.3), font_size=10, color=GREY)
    cols = ["#", "Candidate", "Skills", "Experience", "Culture", "Overall"]
    col_x = [Inches(0.6), Inches(1.1), Inches(3.4), Inches(4.5), Inches(5.6), Inches(6.7)]
    for col, cx in zip(cols, col_x):
        _add_textbox(slide, col, cx, Inches(3.5), Inches(1.1), Inches(0.3),
                     font_size=9, bold=True, color=GOLD)

    for idx, r in enumerate(ranked[:6]):
        y = Inches(3.85) + idx * Inches(0.42)
        row_color = ACCENT if idx % 2 == 0 else DARK
        row_bg = slide.shapes.add_shape(1, Inches(0.6), y - Inches(0.04),
                                         W - Inches(1.2), Inches(0.38))
        row_bg.fill.solid()
        row_bg.fill.fore_color.rgb = row_color
        row_bg.line.fill.background()
        values = [
            str(r.rank),
            r.candidate.name,
            f"{r.evaluation.skills_score:.1f}" if r.evaluation.skills_score else "–",
            f"{r.evaluation.experience_score:.1f}" if r.evaluation.experience_score else "–",
            f"{r.evaluation.culture_score:.1f}" if r.evaluation.culture_score else "–",
            f"{r.score:.1f}",
        ]
        for val, cx in zip(values, col_x):
            c = GOLD if val == str(r.score) else WHITE
            _add_textbox(slide, val, cx, y, Inches(1.1), Inches(0.35), font_size=10, color=c)


def _slide_candidate(prs: Presentation, rc: RankedCandidate, job: Job,
                     interview: Interview | None = None):
    slide = _add_slide(prs)
    _fill_bg(slide, DARK)
    W, H = prs.slide_width, prs.slide_height
    c = rc.candidate
    e = rc.evaluation

    # Header strip
    header = slide.shapes.add_shape(1, 0, 0, W, Inches(1.4))
    header.fill.solid()
    header.fill.fore_color.rgb = ACCENT
    header.line.fill.background()

    _add_textbox(slide, f"#{rc.rank}  {c.name}",
                 Inches(0.5), Inches(0.15), W * 0.6, Inches(0.6),
                 font_size=26, bold=True, color=WHITE)
    _add_textbox(slide, f"{c.current_role or ''}  ·  {c.current_company or ''}",
                 Inches(0.5), Inches(0.75), W * 0.6, Inches(0.4),
                 font_size=13, color=LIGHT)
    _add_textbox(slide, f"{rc.score:.1f}",
                 W - Inches(1.6), Inches(0.1), Inches(1.1), Inches(0.8),
                 font_size=44, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
    _add_textbox(slide, "OVERALL", W - Inches(1.6), Inches(0.9), Inches(1.1), Inches(0.3),
                 font_size=8, color=GREY, align=PP_ALIGN.CENTER)

    # Score bars
    _score_bar(slide, "Skills Match",   e.skills_score or 0,   Inches(0.5), Inches(1.6))
    _score_bar(slide, "Experience",     e.experience_score or 0, Inches(0.5), Inches(2.15))
    _score_bar(slide, "Culture Fit",    e.culture_score or 0,   Inches(0.5), Inches(2.7))
    if e.interview_score:
        _score_bar(slide, "Interview",  e.interview_score,       Inches(0.5), Inches(3.25))

    # Strengths
    _add_textbox(slide, "STRENGTHS", Inches(4.2), Inches(1.5), Inches(3.2), Inches(0.3),
                 font_size=9, bold=True, color=GOLD)
    strengths_text = "\n".join(f"✓  {s}" for s in (e.strengths or [])[:4])
    _add_textbox(slide, strengths_text, Inches(4.2), Inches(1.85), Inches(3.2), Inches(1.2),
                 font_size=10, color=LIGHT)

    # Red flags
    if e.red_flags:
        _add_textbox(slide, "RED FLAGS", Inches(4.2), Inches(3.1), Inches(3.2), Inches(0.3),
                     font_size=9, bold=True, color=GOLD)
        flags_text = "\n".join(f"⚠  {f}" for f in (e.red_flags or [])[:3])
        _add_textbox(slide, flags_text, Inches(4.2), Inches(3.45), Inches(3.2), Inches(0.9),
                     font_size=10, color=LIGHT)

    # Profile details
    details = []
    if c.years_experience:
        details.append(f"Experience: {c.years_experience} yrs")
    if c.location:
        details.append(f"Location: {c.location}")
    if c.availability:
        details.append(f"Available: {c.availability}")
    if c.salary_expectation:
        details.append(f"Salary: {c.salary_expectation:,.0f} {c.salary_currency}")
    _add_textbox(slide, "  ·  ".join(details),
                 Inches(0.5), H - Inches(0.9), W - Inches(1.0), Inches(0.35),
                 font_size=10, color=GREY)

    # Reasoning
    if e.screening_reasoning:
        _add_textbox(slide, e.screening_reasoning[:200] + ("…" if len(e.screening_reasoning) > 200 else ""),
                     Inches(0.5), H - Inches(1.55), W - Inches(1.0), Inches(0.55),
                     font_size=9, color=GREY, italic=True)


def _slide_comparison(prs: Presentation, top3: list[RankedCandidate]):
    slide = _add_slide(prs)
    _fill_bg(slide, DARK)
    W, H = prs.slide_width, prs.slide_height

    _add_textbox(slide, "TOP 3 COMPARISON", Inches(0.5), Inches(0.3),
                 W - Inches(1.0), Inches(0.4), font_size=11, color=GREY)
    _add_textbox(slide, "Head-to-Head", Inches(0.5), Inches(0.7),
                 W - Inches(1.0), Inches(0.6), font_size=28, bold=True, color=WHITE)

    categories = ["Skills Match", "Experience", "Culture Fit"]
    score_getters = [
        lambda e: e.skills_score or 0,
        lambda e: e.experience_score or 0,
        lambda e: e.culture_score or 0,
    ]
    candidate_colors = [GOLD, RGBColor(0x4E, 0xC9, 0xB0), RGBColor(0xC5, 0x86, 0xC0)]

    col_w = Inches(2.2)
    for ci, rc in enumerate(top3[:3]):
        cx = Inches(0.5) + ci * (col_w + Inches(0.2))
        _add_textbox(slide, rc.candidate.name, cx, Inches(1.6), col_w, Inches(0.35),
                     font_size=13, bold=True, color=candidate_colors[ci], align=PP_ALIGN.CENTER)
        _add_textbox(slide, f"Score: {rc.score:.1f}/10", cx, Inches(1.95), col_w, Inches(0.3),
                     font_size=11, color=WHITE, align=PP_ALIGN.CENTER)

        for ri, (cat, getter) in enumerate(zip(categories, score_getters)):
            y = Inches(2.6) + ri * Inches(0.9)
            score = getter(rc.evaluation)
            _add_textbox(slide, cat, cx, y, col_w, Inches(0.28),
                         font_size=9, color=GREY, align=PP_ALIGN.CENTER)
            _add_textbox(slide, f"{score:.1f}", cx, y + Inches(0.28), col_w, Inches(0.5),
                         font_size=28, bold=True, color=candidate_colors[ci], align=PP_ALIGN.CENTER)


def _slide_recommendation(prs: Presentation, ranked: list[RankedCandidate], job: Job):
    slide = _add_slide(prs)
    _fill_bg(slide, DARK)
    W, H = prs.slide_width, prs.slide_height

    _add_textbox(slide, "AI RECOMMENDATION", Inches(0.5), Inches(0.3),
                 W - Inches(1.0), Inches(0.35), font_size=10, color=GREY)
    _add_textbox(slide, "Next Steps", Inches(0.5), Inches(0.65),
                 W - Inches(1.0), Inches(0.6), font_size=28, bold=True, color=WHITE)

    bar = slide.shapes.add_shape(1, Inches(0.5), Inches(1.3), Inches(0.04), Inches(1.0))
    bar.fill.solid()
    bar.fill.fore_color.rgb = GOLD
    bar.line.fill.background()

    top = ranked[0] if ranked else None
    if top:
        _add_textbox(slide, f"Recommended to proceed: {top.candidate.name}",
                     Inches(0.7), Inches(1.3), W - Inches(1.4), Inches(0.45),
                     font_size=16, bold=True, color=WHITE)
        reason = (top.evaluation.screening_reasoning or "")[:300]
        _add_textbox(slide, reason, Inches(0.7), Inches(1.75), W - Inches(1.4), Inches(0.8),
                     font_size=11, color=LIGHT, italic=True)

    steps = [
        "1.  Review shortlisted candidates and approve for interview stage",
        "2.  Schedule structured interviews using the generated question sets",
        "3.  Upload interview notes / transcript for AI evaluation",
        "4.  Final candidate report will be generated after interview stage",
        "5.  All candidate data has been saved to the intelligence database",
    ]
    for i, step in enumerate(steps):
        _add_textbox(slide, step, Inches(0.6), Inches(2.9) + i * Inches(0.5),
                     W - Inches(1.2), Inches(0.4), font_size=12, color=LIGHT)

    _add_textbox(slide, "AI Agency  ·  Hiring Intelligence Platform",
                 Inches(0.5), H - Inches(0.6), W - Inches(1.0), Inches(0.3),
                 font_size=9, color=GREY, align=PP_ALIGN.CENTER)


# ─── Main export ──────────────────────────────────────────────────────────────

def generate_presentation(
    job: Job,
    client_name: str,
    ranked: list[RankedCandidate],
    interviews: dict[str, Interview] | None = None,
    output_dir: str = "output",
) -> str:
    """
    Generate a polished .pptx report for the hiring process.
    Returns path to the saved file.

    The .pptx is directly importable into Prezi (File → Import PowerPoint).
    """
    interviews = interviews or {}
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    _slide_cover(prs, job, client_name)
    _slide_overview(prs, job, ranked)

    for rc in ranked:
        interview = interviews.get(rc.candidate.candidate_id)
        _slide_candidate(prs, rc, job, interview)

    if len(ranked) >= 2:
        _slide_comparison(prs, ranked[:3])

    _slide_recommendation(prs, ranked, job)

    Path(output_dir).mkdir(parents=True, exist_ok=True)
    safe_title = job.title.replace(" ", "_").replace("/", "-")[:40]
    filename = f"{safe_title}_{job.job_id[:8]}.pptx"
    out_path = str(Path(output_dir) / filename)
    prs.save(out_path)

    return out_path
